from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "data" / "curved_real_scale_railway_tunnel_t0t5" / "raycast_vs_regular_comparison"
OUT = OUT_DIR / "step6_short_complete_workflow_report_v2.pptx"
NOTES = OUT_DIR / "step6_short_complete_workflow_report_v2_notes.md"
CHART = OUT_DIR / "step6_error_trend_chart.png"
SECTION_IMAGE = ROOT / "section_annotated.png"

BLUE = RGBColor(15, 76, 129)
BLUE2 = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(239, 246, 255)
RED = RGBColor(220, 38, 38)
ORANGE = RGBColor(245, 158, 11)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
GRAY = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(15, 23, 42)

LABELS = ["T0", "T1", "T2", "T3", "T4", "T5"]
GROUND_TRUTH = [0.0, -10.0, -22.0, -38.0, -58.0, -80.0]
REGULAR = [0.0, -9.9, -21.7, -37.6, -57.3, -79.1]
RAYCAST = [0.0, -10.243, -21.632, -37.070, -56.607, -77.900]
REGULAR_ERR = ["-", "1.00%", "1.36%", "1.05%", "1.21%", "1.13%"]
RAYCAST_ERR = ["-", "2.43%", "1.67%", "2.45%", "2.40%", "2.63%"]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()


def header(slide, title, kicker="STEP 6 VALIDATION"):
    bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    kicker_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.11), Inches(4.5), Inches(0.35))
    p = kicker_box.text_frame.paragraphs[0]
    p.text = kicker
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.85), Inches(12.0), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLACK


def text(slide, x, y, w, h, value, size=15, color=SLATE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    p = frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return box


def bullets(slide, x, y, w, h, items, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)
    return box


def card(slide, x, y, w, h, title, body, fill=LIGHT_BLUE, line=BLUE2, title_size=14, body_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.3)
    text(slide, x + 0.15, y + 0.13, w - 0.3, 0.3, title, title_size, line, True, PP_ALIGN.CENTER)
    text(slide, x + 0.15, y + 0.52, w - 0.3, h - 0.62, body, body_size, SLATE, False, PP_ALIGN.CENTER)


def metric(slide, x, y, label, value, color):
    card(slide, x, y, 2.55, 1.15, label, value, RGBColor(248, 250, 252), color, 11, 24)


def arrow(slide, x, y, w=0.58):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()


def result_table(slide, x, y, w, h):
    data = [["Time", "GT", "Regular", "Raycast", "Reg err", "Ray err"]]
    for index, label in enumerate(LABELS):
        data.append(
            [
                label,
                f"{GROUND_TRUTH[index]:.1f}",
                f"{REGULAR[index]:.3f}".rstrip("0").rstrip("."),
                f"{RAYCAST[index]:.3f}".rstrip("0").rstrip("."),
                REGULAR_ERR[index],
                RAYCAST_ERR[index],
            ]
        )
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col, width in enumerate([0.7, 1.0, 1.1, 1.15, 1.0, 1.0]):
        table.columns[col].width = Inches(width)
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if row_index == 0 else (LIGHT if row_index % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10 if row_index else 9.5)
            p.font.bold = row_index == 0
            p.font.color.rgb = WHITE if row_index == 0 else BLACK
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)


def meaning_table(slide, x, y, w, h):
    data = [
        ["Dataset", "What it proves"],
        ["Ground truth", "Known settlement from Blender; reference for error calculation."],
        ["Regular clean", "Algorithm accuracy under clean, ideal lining data."],
        ["Raycast TLS", "Field-like robustness with scan noise, occlusion and dropout."],
    ]
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(7.7)
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if row_index == 0 else (LIGHT if row_index % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = row_index == 0 or col_index == 0
            p.font.color.rgb = WHITE if row_index == 0 else BLACK
            p.alignment = PP_ALIGN.CENTER if col_index == 0 else PP_ALIGN.LEFT


# 1 Title
slide = prs.slides.add_slide(blank)
bg(slide)
text(slide, 0.8, 0.8, 11.8, 0.8, "Tunnel Crown Settlement Validation Workflow", 34, BLUE, True, PP_ALIGN.CENTER)
text(slide, 1.2, 1.75, 10.8, 0.45, "Blender ground truth → Regular clean → Raycast TLS → Step 6 accuracy check", 18, SLATE, False, PP_ALIGN.CENTER)
metric(slide, 1.25, 3.0, "Regular MAPE", "1.15%", GREEN)
metric(slide, 5.35, 3.0, "Raycast MAPE", "2.315%", ORANGE)
metric(slide, 9.45, 3.0, "Measured point", "Ch 52.0m", BLUE2)
text(slide, 1.4, 5.25, 10.5, 0.5, "Main metric: Crown settlement / Lún đỉnh hầm", 22, RED, True, PP_ALIGN.CENTER)
text(slide, 1.4, 6.05, 10.5, 0.35, "Dataset: curved railway tunnel T0–T5", 14, SLATE, False, PP_ALIGN.CENTER)

# 2 Ground truth goal
slide = prs.slides.add_slide(blank)
header(slide, "Goal: validate the tool using known settlement")
bullets(
    slide,
    0.85,
    1.75,
    5.4,
    3.7,
    [
        "Create one controlled tunnel model in Blender.",
        "Apply known crown settlement from T0 to T5.",
        "Run Step 6 on generated point clouds.",
        "Compare tool output with known ground truth.",
    ],
    18,
)
card(
    slide,
    6.9,
    1.7,
    5.3,
    3.9,
    "Ground truth settlement at Crown / Ch 52.0m",
    "T0 = 0 mm\nT1 = -10 mm\nT2 = -22 mm\nT3 = -38 mm\nT4 = -58 mm\nT5 = -80 mm",
    RGBColor(254, 242, 242),
    RED,
    16,
    18,
)
text(slide, 1.2, 6.1, 10.8, 0.35, "Because the true deformation is known, the measured tool error can be calculated directly.", 14, GREEN, True, PP_ALIGN.CENTER)

# 3 Workflow
slide = prs.slides.add_slide(blank)
header(slide, "Complete workflow: one model, three comparable versions")
card(slide, 0.55, 2.0, 2.0, 1.25, "3D Blender model", "Curved railway tunnel with real-scale components.", LIGHT_BLUE, BLUE2)
arrow(slide, 2.7, 2.48)
card(slide, 3.35, 1.35, 2.05, 1.1, "Ground truth", "Known settlement values.", RGBColor(254, 242, 242), RED)
card(slide, 3.35, 3.0, 2.05, 1.1, "Same crown point", "Crown / Ch 52.0m.", RGBColor(240, 253, 244), GREEN)
arrow(slide, 5.55, 2.48)
card(slide, 6.2, 1.25, 2.05, 1.25, "Regular clean", "Direct clean lining export.", RGBColor(239, 246, 255), BLUE2)
card(slide, 6.2, 3.05, 2.05, 1.25, "Raycast TLS", "Simulated scanner noise and occlusion.", RGBColor(255, 247, 237), ORANGE)
arrow(slide, 8.45, 2.48)
card(slide, 9.1, 2.0, 1.75, 1.25, "Step 6", "Measure crown settlement.", RGBColor(240, 253, 244), GREEN)
arrow(slide, 11.0, 2.48)
card(slide, 11.55, 2.0, 1.2, 1.25, "Error", "mm\n%", LIGHT, RED)
text(slide, 1.0, 5.35, 5.3, 0.45, "Error mm = Tool result − Ground truth", 17, BLACK, True, PP_ALIGN.CENTER)
text(slide, 7.0, 5.35, 5.3, 0.45, "Error % = |Error mm| / |Ground truth| × 100", 17, BLACK, True, PP_ALIGN.CENTER)

# 4 Dataset meaning
slide = prs.slides.add_slide(blank)
header(slide, "What each dataset proves")
meaning_table(slide, 1.75, 1.7, 9.7, 2.7)
text(slide, 1.5, 5.0, 10.4, 0.45, "All three versions come from the same Blender model and are measured at the same crown location.", 16, GREEN, True, PP_ALIGN.CENTER)
text(slide, 1.5, 5.8, 10.4, 0.45, "This separates algorithm error from field-like scan effects.", 16, BLUE, True, PP_ALIGN.CENTER)

# 5 Measured point evidence
slide = prs.slides.add_slide(blank)
header(slide, "Measured point: same crown location for every epoch")
metric(slide, 0.85, 1.55, "Point", "Crown", RED)
metric(slide, 3.75, 1.55, "Location", "Ch 52.0m", BLUE2)
metric(slide, 6.65, 1.55, "Epochs", "T0–T5", GREEN)
metric(slide, 9.55, 1.55, "Display", "2D marker", ORANGE)
if SECTION_IMAGE.exists():
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(3.05), Inches(8.75), Inches(2.55))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = RED
    panel.line.width = Pt(1.2)
    slide.shapes.add_picture(str(SECTION_IMAGE), Inches(2.45), Inches(3.18), width=Inches(8.45), height=Inches(2.25))
else:
    card(slide, 2.3, 3.05, 8.75, 2.55, "2D section marker", "Marker identifies the crown measurement point on the displayed section.", RGBColor(254, 242, 242), RED, 18, 16)
text(slide, 1.1, 6.12, 11.1, 0.45, "The marker is visual evidence only; the measurement values remain true and are not changed by display scaling.", 13.5, SLATE, False, PP_ALIGN.CENTER)

# 6 Pass criteria
slide = prs.slides.add_slide(blank)
header(slide, "Validation criteria: both branches pass")
card(slide, 1.1, 1.75, 4.8, 2.6, "Regular clean criterion", "Pass if MAPE < 2%\n\nActual result: 1.15%\n\nStatus: PASS", RGBColor(240, 253, 244), GREEN, 18, 18)
card(slide, 7.35, 1.75, 4.8, 2.6, "Raycast TLS criterion", "Pass if MAPE < 5%\n\nActual result: 2.315%\n\nStatus: PASS", RGBColor(255, 247, 237), ORANGE, 18, 18)
text(slide, 1.2, 5.25, 10.9, 0.5, "Tool passes controlled validation: clean accuracy is high and field-like robustness is acceptable.", 18, BLUE, True, PP_ALIGN.CENTER)
text(slide, 1.2, 6.1, 10.9, 0.35, "Raycast error is expected to be higher because it includes scanner noise, occlusion and missing points.", 13.5, SLATE, False, PP_ALIGN.CENTER)

# 7 Results
slide = prs.slides.add_slide(blank)
header(slide, "Accuracy result: tool output vs ground truth")
result_table(slide, 0.65, 1.45, 6.0, 4.25)
metric(slide, 7.4, 1.65, "Regular MAPE", "1.15%", GREEN)
metric(slide, 10.15, 1.65, "Raycast MAPE", "2.315%", ORANGE)
text(slide, 7.45, 3.35, 4.75, 1.2, "Regular is more accurate because the surface is clean. Raycast error is higher because it includes field-like scan effects.", 15, SLATE)
text(slide, 7.45, 5.0, 4.75, 0.95, "Both still follow the same T0→T5 settlement trend, so Step 6 quantifies crown settlement reliably.", 15, GREEN, True)
text(slide, 0.8, 6.35, 11.8, 0.35, "T0 is not used in MAPE because its ground truth is 0 mm.", 11, SLATE, False, PP_ALIGN.CENTER)

# 8 Conclusion
slide = prs.slides.add_slide(blank)
header(slide, "Conclusion")
text(slide, 1.0, 1.45, 11.2, 0.55, "The workflow validates Step 6 against a known Blender ground truth.", 24, BLUE, True, PP_ALIGN.CENTER)
metric(slide, 1.4, 2.6, "Clean data error", "1.15%", GREEN)
metric(slide, 5.35, 2.6, "Field-like error", "2.315%", ORANGE)
metric(slide, 9.3, 2.6, "Measured at", "Crown / Ch 52m", BLUE2)
bullets(
    slide,
    1.7,
    4.35,
    9.9,
    1.35,
    [
        "Regular clean validates the algorithm under ideal conditions.",
        "Raycast TLS validates robustness under field-like scanning conditions.",
        "Next step: test with real TLS scans and calibrate preprocessing.",
    ],
    18,
)
text(slide, 1.1, 6.25, 11.1, 0.35, "Key message: Step 6 measures crown settlement accurately enough for controlled validation.", 15, RED, True, PP_ALIGN.CENTER)

prs.save(OUT)
NOTES.write_text(
    """# Short presentation script v2

1. I created one controlled curved railway tunnel in Blender.
2. I applied known crown settlement values from T0 to T5, so this is the ground truth.
3. From the same model, I generated two datasets: regular clean and raycast field-like TLS.
4. Regular clean tests ideal algorithm accuracy; raycast TLS tests field-like robustness.
5. Step 6 measures only crown settlement at the same location: Crown / Ch 52.0m.
6. I compare tool output with ground truth using error in mm and percent.
7. Validation criteria are Regular MAPE < 2% and Raycast MAPE < 5%.
8. Results pass: Regular MAPE = 1.15%, Raycast MAPE = 2.315%.
9. The raycast error is higher because it includes noise and occlusion, but the settlement trend remains correct from T0 to T5.
""",
    encoding="utf-8",
)
print(OUT)
print(NOTES)
