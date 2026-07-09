from pathlib import Path
import math
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / 'data' / 'curved_real_scale_railway_tunnel_t0t5' / 'raycast_vs_regular_comparison'
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX = OUT_DIR / 'step6_crown_settlement_workflow_professor_report.pptx'
CHART = OUT_DIR / 'step6_error_trend_chart.png'

BLUE = RGBColor(15, 76, 129)
BLUE2 = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(239, 246, 255)
RED = RGBColor(220, 38, 38)
ORANGE = RGBColor(245, 158, 11)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
GRAY = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(15, 23, 42)

# Data
labels = ['T0', 'T1', 'T2', 'T3', 'T4', 'T5']
gt = [0.0, -10.0, -22.0, -38.0, -58.0, -80.0]
regular = [0.0, -9.9, -21.7, -37.6, -57.3, -79.1]
raycast = [0.0, -10.243, -21.632, -37.070, -56.607, -77.900]
regular_err = [0.0, 1.0, 1.364, 1.053, 1.207, 1.125]
raycast_err = [0.0, 2.43, 1.673, 2.447, 2.402, 2.625]

plt.figure(figsize=(7.2, 3.6), dpi=180)
plt.plot(labels, regular_err, marker='o', linewidth=2.5, color='#2563EB', label='Regular clean')
plt.plot(labels, raycast_err, marker='o', linewidth=2.5, color='#F59E0B', label='Raycast field-like')
plt.axhspan(0, 2, color='#DCFCE7', alpha=0.35, label='Target range')
plt.ylim(0, 3.2)
plt.ylabel('Absolute error (%)')
plt.xlabel('Time epoch')
plt.title('Crown Settlement Error Trend')
plt.grid(True, alpha=0.25)
plt.legend(loc='upper left', frameon=True)
plt.tight_layout()
plt.savefig(CHART, transparent=False, facecolor='white')
plt.close()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()


def add_header(slide, title, kicker='STEP 6 VALIDATION'):
    add_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.11), Inches(3.0), Inches(0.35))
    p = tx.text_frame.paragraphs[0]
    p.text = kicker; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHITE
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(11.8), Inches(0.55))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = BLACK
    return slide


def add_text(slide, x, y, w, h, text, size=16, color=SLATE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    if align: p.alignment = align
    return box


def add_bullets(slide, x, y, w, h, items, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.size = Pt(size); p.font.color.rgb = SLATE
        p.space_after = Pt(8)
    return box


def add_card(slide, x, y, w, h, title, body, color=LIGHT_BLUE, accent=BLUE2):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.color.rgb = accent; rect.line.width = Pt(1.2)
    add_text(slide, x+0.18, y+0.15, w-0.36, 0.3, title, size=14, color=accent, bold=True)
    add_text(slide, x+0.18, y+0.55, w-0.36, h-0.65, body, size=12, color=SLATE)
    return rect


def add_metric(slide, x, y, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1.1))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(248,250,252)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    add_text(slide, x+0.15, y+0.12, 2.2, 0.28, label, size=11, color=SLATE, bold=True)
    add_text(slide, x+0.15, y+0.44, 2.2, 0.5, value, size=24, color=color, bold=True)


def add_table(slide, x, y, w, h, data, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(9 if rows > 6 else 11)
                p.font.color.rgb = WHITE if r == 0 else BLACK
                p.font.bold = r == 0
            cell.fill.solid(); cell.fill.fore_color.rgb = header_fill if r == 0 else (RGBColor(248,250,252) if r % 2 else WHITE)
    return tbl

# 1
s = prs.slides.add_slide(blank); add_bg(s)
add_text(s, 0.8, 0.75, 11.8, 0.8, 'Tunnel Crown Settlement Monitoring Workflow', 34, BLACK, True)
add_text(s, 0.82, 1.55, 10.8, 0.45, 'Blender Ground Truth → Raycast TLS → Step 6 Validation', 20, BLUE2, True)
add_card(s, 0.9, 2.45, 3.5, 1.4, 'Dataset', 'curved_real_scale_railway_tunnel_t0t5\nT0–T5 railway tunnel settlement', LIGHT_BLUE, BLUE2)
add_card(s, 4.9, 2.45, 3.5, 1.4, 'Measured Point', 'Crown / Đỉnh hầm\nLocation: Ch 52.0m', RGBColor(254,242,242), RED)
add_card(s, 8.9, 2.45, 3.5, 1.4, 'Main Result', 'Regular MAPE 1.15%\nRaycast MAPE 2.315%', RGBColor(240,253,244), GREEN)
add_text(s, 0.9, 6.65, 11.5, 0.35, 'CBNU Smart Structure Lab · Step 6 Crown-First Validation', 12, SLATE)

# 2
s = prs.slides.add_slide(blank); add_header(s, 'Research Goal')
add_bullets(s, 0.8, 1.7, 5.5, 3.4, [
    'Monitor tunnel settlement over time using point-cloud scan epochs.',
    'Measure the same engineering location: Crown / Đỉnh hầm at Ch 52.0m.',
    'Report total settlement, new movement from previous epoch, and status.',
    'Keep the main Step 6 decision metric simple: crown_settlement_mm.'
], 17)
add_card(s, 7.0, 1.6, 4.9, 1.2, 'Decision Output', 'OK / Warning / Danger based on absolute crown settlement.', RGBColor(255,247,237), ORANGE)
add_card(s, 7.0, 3.1, 4.9, 1.2, 'Warning Threshold', 'Warning ≥ 10 mm', RGBColor(254,249,195), ORANGE)
add_card(s, 7.0, 4.6, 4.9, 1.2, 'Danger Threshold', 'Danger ≥ 25 mm', RGBColor(254,226,226), RED)

# 3 workflow
s = prs.slides.add_slide(blank); add_header(s, 'Overall Workflow')
steps = ['Blender tunnel\nmodel', 'Ground truth\nT0–T5', 'Regular clean\ndataset', 'Raycast TLS\nfield-like', 'Step 6\nanalysis', 'Accuracy\ncomparison']
colors = [BLUE2, RED, GREEN, ORANGE, BLUE, SLATE]
for i, step in enumerate(steps):
    x = 0.55 + i*2.08
    add_card(s, x, 2.55, 1.65, 1.25, f'{i+1}', step, RGBColor(248,250,252), colors[i])
    if i < len(steps)-1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+1.62), Inches(2.95), Inches(0.55), Inches(0.35))
        arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()
add_text(s, 1.0, 5.0, 11.0, 0.5, 'The same crown location is tracked across all epochs and both data branches.', 18, BLUE, True, PP_ALIGN.CENTER)

# 4
s = prs.slides.add_slide(blank); add_header(s, 'Blender Ground Truth Setup')
add_bullets(s, 0.8, 1.55, 5.0, 2.8, [
    'Curved railway tunnel generated at real scale.',
    'Known crown settlement is injected by epoch.',
    'Ground truth gives a controlled benchmark for Step 6.'
], 17)
table = [['Epoch','T0','T1','T2','T3','T4','T5'], ['GT mm','0','-10','-22','-38','-58','-80']]
add_table(s, 0.8, 4.65, 6.1, 0.9, table)
# simple tunnel sketch
arc = s.shapes.add_shape(MSO_SHAPE.ARC, Inches(7.45), Inches(2.0), Inches(3.8), Inches(2.2))
arc.line.color.rgb = BLUE2; arc.line.width = Pt(4)
for i, val in enumerate(gt[1:]):
    add_text(s, 7.2+i*0.85, 4.35+0.08*i, 0.9, 0.3, f'T{i+1}\n{val}mm', 10, RED, True, PP_ALIGN.CENTER)
add_text(s, 7.2, 5.55, 4.6, 0.45, 'Known truth for validation', 18, RED, True, PP_ALIGN.CENTER)

# 5
s = prs.slides.add_slide(blank); add_header(s, 'Two Test Data Branches')
add_card(s, 1.0, 1.7, 5.2, 3.8, 'Regular clean', 'Clean lining surface\nLow noise and no TLS occlusion\nPurpose: ideal algorithm check\nExpected: lowest error', RGBColor(240,253,244), GREEN)
add_card(s, 7.1, 1.7, 5.2, 3.8, 'Raycast field-like TLS', 'Simulated scanning stations\nNoise, dropout, and occlusion\nPurpose: field readiness check\nExpected: higher but acceptable error', RGBColor(255,247,237), ORANGE)
add_text(s, 1.2, 6.1, 10.8, 0.45, 'Regular tests the best-case algorithm; raycast tests practical field behavior.', 18, BLUE, True, PP_ALIGN.CENTER)

# 6
s = prs.slides.add_slide(blank); add_header(s, 'Step 6 Crown-First Design')
add_metric(s, 0.9, 1.6, 'Main metric', 'crown_settlement_mm', BLUE2)
add_metric(s, 3.9, 1.6, 'Measured point', 'Crown / Đỉnh hầm', RED)
add_metric(s, 6.9, 1.6, 'Location', 'Ch 52.0m', GREEN)
add_metric(s, 9.9, 1.6, 'Rule', 'abs(crown)', ORANGE)
add_bullets(s, 1.0, 3.4, 10.8, 2.2, [
    'M3C2 and p95 remain supplementary context, not the main settlement result.',
    'Warning starts at 10 mm; Danger starts at 25 mm.',
    'Chart, table, 2D marker, CSV, Excel, and PDF should all point to the same metric.'
], 17)

# 7
s = prs.slides.add_slide(blank); add_header(s, 'Step 6 UI Improvements')
ui_table = [['Time','Location','Measured point','Crown','New move','Result'], ['T5','Ch 52.0m','Crown / Đỉnh hầm','-79.1','-21.8','Danger']]
add_table(s, 0.8, 1.7, 7.5, 1.0, ui_table)
add_bullets(s, 0.95, 3.2, 6.6, 2.4, [
    'Measured point markers for T1–T5 on the 2D section.',
    'Checkboxes hide/show marker per epoch; lines remain visible.',
    'Crown visual x10 magnifies display only; measurements stay true.'
], 16)
add_card(s, 8.6, 1.7, 3.7, 3.9, 'User-facing goal', 'The engineer can answer three questions quickly:\n\n1. Where was it measured?\n2. How much did it settle?\n3. Is it OK, Warning, or Danger?', RGBColor(239,246,255), BLUE2)

# 8 table
s = prs.slides.add_slide(blank); add_header(s, 'Accuracy Result Table')
data = [['Time','GT','Regular','Raycast','Reg err %','Ray err %']]
for i in range(len(labels)):
    data.append([labels[i], f'{gt[i]:.0f}', f'{regular[i]:.3g}', f'{raycast[i]:.3f}'.rstrip('0').rstrip('.'), f'{regular_err[i]:.2f}', f'{raycast_err[i]:.2f}'])
add_table(s, 0.7, 1.45, 8.1, 3.2, data)
add_metric(s, 9.25, 1.55, 'Regular MAPE', '1.15%', GREEN)
add_metric(s, 9.25, 3.0, 'Raycast MAPE', '2.315%', ORANGE)
add_text(s, 0.9, 5.55, 11.4, 0.55, 'Raycast is harder because it includes field-like noise and occlusion, but error remains around 2–3%.', 18, BLUE, True, PP_ALIGN.CENTER)

# 9 chart
s = prs.slides.add_slide(blank); add_header(s, 'Error Trend Visualization')
s.shapes.add_picture(str(CHART), Inches(0.9), Inches(1.45), width=Inches(7.6))
add_card(s, 8.9, 1.65, 3.3, 1.3, 'Interpretation', 'Regular clean stays close to ~1% error.', RGBColor(240,253,244), GREEN)
add_card(s, 8.9, 3.25, 3.3, 1.3, 'Field-like result', 'Raycast error is higher but stable and acceptable.', RGBColor(255,247,237), ORANGE)
add_card(s, 8.9, 4.85, 3.3, 1.3, 'Trend', 'Settlement magnitude follows T0→T5 correctly.', LIGHT_BLUE, BLUE2)

# 10 2D evidence mock
s = prs.slides.add_slide(blank); add_header(s, '2D Section Visual Evidence')
# draw simplified multi-lines
x0, y0, w, h = 1.0, 1.5, 8.2, 4.7
plot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(y0), Inches(w), Inches(h))
plot.fill.solid(); plot.fill.fore_color.rgb = WHITE; plot.line.color.rgb = GRAY
for i, lab in enumerate(labels):
    y = y0 + 1.0 + i*0.18
    color = [RGBColor(59,130,246), RGBColor(16,185,129), RGBColor(101,163,13), RGBColor(245,158,11), RGBColor(239,68,68), RGBColor(220,38,38)][i]
    pts = [(x0+0.5,y+1.1),(x0+2.0,y+0.4),(x0+4.2,y+0.25),(x0+6.1,y+0.35),(x0+7.6,y+1.15)]
    for a,b in zip(pts[:-1], pts[1:]):
        line = s.shapes.add_connector(1, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1]))
        line.line.color.rgb = color; line.line.width = Pt(2)
    if i > 0:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0+4.05+0.09*i), Inches(y+0.18), Inches(0.14), Inches(0.14))
        circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.color.rgb = WHITE
        add_text(s, x0+4.15+0.09*i, y+0.02, 1.1, 0.22, f'{lab}: {regular[i]:.1f}mm', 8, RED, True)
add_text(s, 1.1, 6.35, 7.8, 0.35, 'Measured points are visual markers on the displayed outline; measurements remain true.', 13, SLATE)
add_card(s, 9.7, 1.8, 2.6, 3.3, 'QA rule', 'Marker must sit on the matching T line.\n\nCheckboxes hide markers only, not the section lines.', RGBColor(254,242,242), RED)

# 11
s = prs.slides.add_slide(blank); add_header(s, 'Field Readiness')
add_bullets(s, 0.9, 1.5, 6.1, 3.8, [
    'Curved tunnel centerline supports non-straight railway tunnel geometry.',
    'Raycast branch simulates multiple TLS stations and field occlusion.',
    'Robust 2D outline reduces visual spikes from rails, pipes, and cables.',
    'Simple Step 6 table is readable for field engineers.'
], 16)
add_card(s, 7.6, 1.7, 4.5, 2.0, 'Practical limitation', 'For raw field scans, denoise and lining extraction should be performed before official Step 6 evaluation.', RGBColor(255,247,237), ORANGE)
add_card(s, 7.6, 4.1, 4.5, 1.5, 'What is validated?', 'Crown settlement trend and magnitude are recovered accurately in both clean and raycast conditions.', RGBColor(240,253,244), GREEN)

# 12
s = prs.slides.add_slide(blank); add_header(s, 'Conclusion')
add_metric(s, 0.9, 1.55, 'Regular MAPE', '1.15%', GREEN)
add_metric(s, 3.9, 1.55, 'Raycast MAPE', '2.315%', ORANGE)
add_metric(s, 6.9, 1.55, 'Measured at', 'Ch 52.0m', BLUE2)
add_bullets(s, 1.0, 3.25, 10.8, 2.3, [
    'The tool measures crown settlement close to known Blender ground truth.',
    'Step 6 now uses one main metric consistently: crown_settlement_mm.',
    'The workflow is ready for additional real TLS scan validation.'
], 18)
add_text(s, 1.0, 6.25, 10.8, 0.5, 'Next work: test real field scans and calibrate preprocessing for different tunnel types.', 16, BLUE, True, PP_ALIGN.CENTER)

prs.save(PPTX)
print(PPTX)
