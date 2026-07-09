from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / 'data' / 'curved_real_scale_railway_tunnel_t0t5' / 'raycast_vs_regular_comparison'
OUT = OUT_DIR / 'step6_easy_workflow_explanation.pptx'
NOTES = OUT_DIR / 'step6_easy_workflow_explanation_notes.md'
SECTION_IMAGE = ROOT / 'section_annotated.png'
CHART = OUT_DIR / 'step6_error_trend_chart.png'

BLUE = RGBColor(15, 76, 129)
BLUE2 = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(239, 246, 255)
RED = RGBColor(220, 38, 38)
ORANGE = RGBColor(245, 158, 11)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
GRAY = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(15, 23, 42)

LABELS = ['T0', 'T1', 'T2', 'T3', 'T4', 'T5']
GT = [0.0, -10.0, -22.0, -38.0, -58.0, -80.0]
REG = [0.0, -9.9, -21.7, -37.6, -57.3, -79.1]
RAY = [0.0, -10.243, -21.632, -37.070, -56.607, -77.900]
REG_ERR = ['-', '1.00%', '1.36%', '1.05%', '1.21%', '1.13%']
RAY_ERR = ['-', '2.43%', '1.67%', '2.45%', '2.40%', '2.63%']

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

def text(slide, x, y, w, h, value, size=15, color=SLATE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = value; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    if align:
        p.alignment = align
    return box

def header(slide, title, step=''):
    bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    text(slide, 0.45, 0.12, 3.2, 0.35, step or 'STEP 6 WORKFLOW', 10, WHITE, True)
    text(slide, 0.7, 0.86, 12.0, 0.62, title, 27, BLACK, True)

def card(slide, x, y, w, h, title, body, fill=LIGHT_BLUE, line=BLUE2, title_size=15, body_size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    text(slide, x+0.16, y+0.14, w-0.32, 0.34, title, title_size, line, True, PP_ALIGN.CENTER)
    text(slide, x+0.18, y+0.58, w-0.36, h-0.68, body, body_size, SLATE, False, PP_ALIGN.CENTER)

def metric(slide, x, y, label, value, color):
    card(slide, x, y, 2.55, 1.1, label, value, LIGHT, color, 11, 23)

def bullets(slide, x, y, w, h, items, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = SLATE
        p.space_after = Pt(9)
    return box

def arrow(slide, x, y, w=0.55):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shp.fill.solid(); shp.fill.fore_color.rgb = GRAY
    shp.line.fill.background()

def result_table(slide, x, y, w, h):
    rows = [['Time','Ground truth','Regular tool','Raycast tool','Regular error','Raycast error']]
    for i, label in enumerate(LABELS):
        rows.append([label, f'{GT[i]:.1f}', f'{REG[i]:.3f}'.rstrip('0').rstrip('.'), f'{RAY[i]:.3f}'.rstrip('0').rstrip('.'), REG_ERR[i], RAY_ERR[i]])
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [0.72, 1.28, 1.25, 1.25, 1.15, 1.15]
    for c, width in enumerate(widths):
        tbl.columns[c].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r,c); cell.text = value
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE if r == 0 else (LIGHT if r % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9 if r == 0 else 9.5)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else BLACK
            cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)

# 1
s = prs.slides.add_slide(blank); bg(s)
text(s, 0.75, 0.75, 11.9, 0.8, 'Easy Workflow to Validate Tunnel Settlement Tool', 32, BLUE, True, PP_ALIGN.CENTER)
text(s, 1.25, 1.7, 10.9, 0.48, 'From Blender 3D model to ground truth, raycast data, regular data, and tool error comparison', 17, SLATE, False, PP_ALIGN.CENTER)
metric(s, 1.1, 3.0, 'Main point', 'Crown', RED)
metric(s, 4.0, 3.0, 'Location', 'Ch 52.0m', BLUE2)
metric(s, 6.9, 3.0, 'Regular MAPE', '1.15%', GREEN)
metric(s, 9.8, 3.0, 'Raycast MAPE', '2.315%', ORANGE)
text(s, 1.2, 5.25, 10.9, 0.5, 'Goal: prove how much Step 6 differs from known ground truth.', 21, BLACK, True, PP_ALIGN.CENTER)

# 2
s = prs.slides.add_slide(blank); header(s, 'Step 1 — Create the 3D tunnel model in Blender', 'STEP 1')
card(s, 0.9, 1.75, 5.5, 3.6, 'Input model', 'A curved railway tunnel is built at real scale.\n\nIt includes lining, rail, walkway, pipes, cables, lights and equipment.', LIGHT_BLUE, BLUE2, 18, 16)
card(s, 7.0, 1.75, 5.1, 3.6, 'Why this matters', 'The Blender model is controlled by us.\n\nThat means we know exactly what deformation is applied later.', RGBColor(240,253,244), GREEN, 18, 16)
text(s, 1.1, 6.05, 11.1, 0.45, 'Simple meaning: Blender is the source model used to generate all test data.', 17, BLUE, True, PP_ALIGN.CENTER)

# 3
s = prs.slides.add_slide(blank); header(s, 'Step 2 — Create ground truth settlement', 'STEP 2')
bullets(s, 0.95, 1.75, 5.3, 3.5, [
    'Use the same tunnel for T0, T1, T2, T3, T4, T5.',
    'Apply known settlement at the crown.',
    'This known value becomes the answer key.'
], 18)
card(s, 6.85, 1.65, 5.2, 3.95, 'Ground truth at Crown / Ch 52.0m', 'T0 = 0 mm\nT1 = -10 mm\nT2 = -22 mm\nT3 = -38 mm\nT4 = -58 mm\nT5 = -80 mm', RGBColor(254,242,242), RED, 17, 18)
text(s, 1.1, 6.15, 11.1, 0.35, 'Simple meaning: we already know the correct settlement, so we can check the tool.', 16, GREEN, True, PP_ALIGN.CENTER)

# 4
s = prs.slides.add_slide(blank); header(s, 'Step 3 — Create regular clean data without raycasting', 'STEP 3')
card(s, 1.0, 1.7, 5.15, 3.65, 'Regular clean', 'Export the clean tunnel lining surface directly from Blender.\n\nThis data has very little noise and no scanner occlusion.', RGBColor(239,246,255), BLUE2, 19, 16)
card(s, 7.1, 1.7, 5.15, 3.65, 'Purpose', 'Check the measurement algorithm in ideal conditions.\n\nIf this branch is wrong, the tool logic may be wrong.', RGBColor(240,253,244), GREEN, 19, 16)
text(s, 1.1, 6.05, 11.1, 0.45, 'Simple meaning: regular clean tests whether the tool can measure correctly on clean data.', 16, BLUE, True, PP_ALIGN.CENTER)

# 5
s = prs.slides.add_slide(blank); header(s, 'Step 4 — Create raycast TLS data like field scanning', 'STEP 4')
card(s, 0.85, 1.65, 5.25, 3.75, 'Raycast field-like TLS', 'Simulate laser scanning from multiple stations along the curved tunnel.\n\nRays hit the model and create a point cloud.', RGBColor(255,247,237), ORANGE, 19, 16)
card(s, 7.0, 1.65, 5.35, 3.75, 'Field-like effects', 'This data can include noise, missing points, occlusion and objects inside the tunnel.\n\nIt is closer to real TLS data.', RGBColor(254,242,242), RED, 19, 16)
text(s, 1.1, 6.05, 11.1, 0.45, 'Simple meaning: raycast tests whether the tool still works when data looks more realistic.', 16, ORANGE, True, PP_ALIGN.CENTER)

# 6
s = prs.slides.add_slide(blank); header(s, 'Step 5 — Run Step 6 on both datasets', 'STEP 5')
metric(s, 0.8, 1.65, 'Metric', 'Crown', RED)
metric(s, 3.75, 1.65, 'Location', 'Ch 52.0m', BLUE2)
metric(s, 6.7, 1.65, 'Epochs', 'T0–T5', GREEN)
metric(s, 9.65, 1.65, 'Output', 'Result', ORANGE)
bullets(s, 1.0, 3.35, 5.4, 2.05, [
    'Crown settlement = total settlement compared with T0.',
    'New crown move = movement added from the previous epoch.',
    'Result = OK, Warning or Danger.'
], 17)
if SECTION_IMAGE.exists():
    s.shapes.add_picture(str(SECTION_IMAGE), Inches(7.0), Inches(3.15), width=Inches(4.9), height=Inches(2.15))
else:
    card(s, 7.0, 3.15, 4.9, 2.15, '2D evidence', 'Measured point marker shows where the crown is measured.', LIGHT, RED, 18, 15)
text(s, 1.1, 6.1, 11.1, 0.35, 'M3C2/p95 are supplementary checks, not the main conclusion.', 14, SLATE, False, PP_ALIGN.CENTER)

# 7
s = prs.slides.add_slide(blank); header(s, 'Step 6 — Compare three versions to find tool error', 'STEP 6')
card(s, 0.55, 2.0, 2.0, 1.25, 'Ground truth', 'Known answer from Blender.', RGBColor(254,242,242), RED)
arrow(s, 2.7, 2.48)
card(s, 3.35, 1.25, 2.15, 1.25, 'Regular result', 'Tool result on clean data.', RGBColor(239,246,255), BLUE2)
card(s, 3.35, 3.05, 2.15, 1.25, 'Raycast result', 'Tool result on field-like data.', RGBColor(255,247,237), ORANGE)
arrow(s, 5.75, 2.48)
card(s, 6.4, 2.0, 2.3, 1.25, 'Compare', 'Tool − ground truth.', RGBColor(240,253,244), GREEN)
arrow(s, 8.95, 2.48)
card(s, 9.6, 2.0, 2.7, 1.25, 'Final error', 'Error in mm and %.', LIGHT, RED)
text(s, 1.15, 5.25, 5.2, 0.45, 'Error mm = Tool result − Ground truth', 17, BLACK, True, PP_ALIGN.CENTER)
text(s, 6.95, 5.25, 5.2, 0.45, 'Error % = |Error mm| / |Ground truth| × 100', 17, BLACK, True, PP_ALIGN.CENTER)
text(s, 1.1, 6.15, 11.1, 0.35, 'T0 is not used in MAPE because its ground truth is 0 mm.', 13, SLATE, False, PP_ALIGN.CENTER)

# 8
s = prs.slides.add_slide(blank); header(s, 'Accuracy result and conclusion', 'RESULT')
result_table(s, 0.55, 1.45, 6.9, 4.25)
metric(s, 8.05, 1.55, 'Regular MAPE', '1.15%', GREEN)
metric(s, 10.75, 1.55, 'Pass limit', '< 2%', GREEN)
metric(s, 8.05, 3.1, 'Raycast MAPE', '2.315%', ORANGE)
metric(s, 10.75, 3.1, 'Pass limit', '< 5%', ORANGE)
text(s, 8.1, 4.95, 4.6, 0.8, 'Both branches pass the validation criteria.', 18, GREEN, True, PP_ALIGN.CENTER)
text(s, 8.1, 5.9, 4.6, 0.55, 'Raycast error is higher, but acceptable because it simulates field noise and occlusion.', 13, SLATE, False, PP_ALIGN.CENTER)

# 9
s = prs.slides.add_slide(blank); header(s, 'Final message for the professor', 'SUMMARY')
text(s, 0.9, 1.55, 11.6, 0.6, 'The tool is validated by comparing Step 6 results against known Blender ground truth.', 24, BLUE, True, PP_ALIGN.CENTER)
card(s, 1.15, 2.75, 3.25, 1.45, 'Clean branch', 'Algorithm check\nMAPE = 1.15%\nPASS', RGBColor(240,253,244), GREEN, 17, 16)
card(s, 5.05, 2.75, 3.25, 1.45, 'Raycast branch', 'Field-like check\nMAPE = 2.315%\nPASS', RGBColor(255,247,237), ORANGE, 17, 16)
card(s, 8.95, 2.75, 3.25, 1.45, 'Main metric', 'Crown settlement\nCrown / Ch 52.0m', RGBColor(239,246,255), BLUE2, 17, 16)
bullets(s, 1.6, 5.15, 10.2, 1.0, [
    'Regular clean proves the tool works in ideal conditions.',
    'Raycast TLS proves the tool remains accurate with realistic scanning effects.',
    'Next step: test and calibrate with real tunnel TLS scans.'
], 17)

prs.save(OUT)
NOTES.write_text('''# Easy presentation script

1. First, I create a curved railway tunnel model in Blender. This model is controlled, so I can use it as the source for all test data.
2. Then I apply known crown settlement values from T0 to T5. This is the ground truth, or the answer key.
3. Next, I create regular clean data by exporting the clean lining surface directly from Blender. This checks the algorithm under ideal conditions.
4. I also create raycast TLS data by simulating laser scanning from multiple stations. This data includes field-like effects such as noise, occlusion and missing points.
5. I run Step 6 on both datasets. Step 6 measures crown settlement at the same location: Crown / Ch 52.0m.
6. Finally, I compare three versions: ground truth, regular tool result and raycast tool result. The error is calculated in millimeters and percent.
7. The result is Regular MAPE = 1.15% and Raycast MAPE = 2.315%. Both pass the validation criteria, so the tool is validated under controlled conditions.
''', encoding='utf-8')
print(OUT)
print(NOTES)
