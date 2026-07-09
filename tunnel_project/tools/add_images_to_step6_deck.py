from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / 'data' / 'curved_real_scale_railway_tunnel_t0t5' / 'raycast_vs_regular_comparison' / 'step6_crown_settlement_workflow_professor_report.pptx'
OUT = ROOT / 'data' / 'curved_real_scale_railway_tunnel_t0t5' / 'raycast_vs_regular_comparison' / 'step6_crown_settlement_workflow_professor_report_with_images.pptx'
imgs = {
    'track': ROOT / 'docs' / 'img' / 'tunnel_track.png',
    'interior': ROOT / 'docs' / 'img' / 'tunnel_interior.png',
    'ground': ROOT / 'data' / 'tunnel_t0t5_blend' / 'ground_truth_explained_en.png',
    'section': ROOT / 'section_annotated.png',
    'info': ROOT / 'output' / 'section_parameter_infographic.png',
}
prs = Presentation(PPTX)

BLUE = RGBColor(15, 76, 129)
WHITE = RGBColor(255,255,255)
BLACK = RGBColor(15,23,42)
SLATE = RGBColor(51,65,85)
RED = RGBColor(220,38,38)
ORANGE = RGBColor(245,158,11)
GREEN = RGBColor(22,163,74)


def add_caption(slide, x, y, w, text, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.italic = True
    return box


def add_panel(slide, x, y, w, h, fill=WHITE, line=RGBColor(226,232,240)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp

# Slide 1: add tunnel hero image on right with white border
s = prs.slides[0]
add_panel(s, 7.35, 4.15, 4.85, 1.95, WHITE, BLUE)
s.shapes.add_picture(str(imgs['track']), Inches(7.48), Inches(4.25), width=Inches(4.58), height=Inches(1.72))
add_caption(s, 7.48, 6.02, 4.58, 'Representative railway tunnel environment', BLUE)

# Slide 3: add small visual proof image below workflow
s = prs.slides[2]
add_panel(s, 4.25, 5.45, 4.8, 1.12, WHITE, BLUE)
s.shapes.add_picture(str(imgs['interior']), Inches(4.35), Inches(5.52), width=Inches(4.6), height=Inches(0.92))
add_caption(s, 4.35, 6.43, 4.6, 'Workflow targets field-like tunnel scanning conditions', BLUE)

# Slide 4: replace sketch area with ground-truth image crop-like fit
s = prs.slides[3]
add_panel(s, 7.0, 1.55, 5.35, 4.45, WHITE, RED)
s.shapes.add_picture(str(imgs['ground']), Inches(7.18), Inches(1.72), width=Inches(4.95), height=Inches(3.92))
add_caption(s, 7.18, 5.72, 4.95, 'Ground-truth settlement explanation used for validation', RED)

# Slide 6: add section parameter infographic as visual context
s = prs.slides[5]
add_panel(s, 1.0, 5.15, 11.2, 1.35, WHITE, BLUE)
s.shapes.add_picture(str(imgs['info']), Inches(1.15), Inches(5.25), width=Inches(10.9), height=Inches(1.05))
add_caption(s, 1.15, 6.32, 10.9, 'Section parameters are available, but Step 6 conclusion remains crown-first.', BLUE)

# Slide 7: add UI/section image thumbnail
s = prs.slides[6]
add_panel(s, 8.55, 5.15, 3.8, 1.25, WHITE, GREEN)
s.shapes.add_picture(str(imgs['section']), Inches(8.65), Inches(5.22), width=Inches(3.6), height=Inches(1.02))
add_caption(s, 8.65, 6.25, 3.6, '2D section view: marker/outline evidence', GREEN)

# Slide 10: replace mock evidence emphasis with actual section image large
s = prs.slides[9]
add_panel(s, 0.9, 1.35, 8.7, 5.2, WHITE, RED)
s.shapes.add_picture(str(imgs['section']), Inches(1.05), Inches(1.48), width=Inches(8.4), height=Inches(4.76))
add_caption(s, 1.05, 6.28, 8.4, 'Actual 2D section evidence: measured marker should sit on the displayed outline.', RED)

# Slide 11: add field image
s = prs.slides[10]
add_panel(s, 7.45, 5.45, 4.75, 1.35, WHITE, ORANGE)
s.shapes.add_picture(str(imgs['track']), Inches(7.55), Inches(5.52), width=Inches(4.55), height=Inches(1.05))
add_caption(s, 7.55, 6.6, 4.55, 'Field readiness target: railway tunnel TLS workflow', ORANGE)

# Slide 12: add small visual close
s = prs.slides[11]
add_panel(s, 9.65, 5.45, 2.55, 1.15, WHITE, GREEN)
s.shapes.add_picture(str(imgs['interior']), Inches(9.75), Inches(5.52), width=Inches(2.35), height=Inches(0.88))

prs.save(OUT)
print(OUT)
