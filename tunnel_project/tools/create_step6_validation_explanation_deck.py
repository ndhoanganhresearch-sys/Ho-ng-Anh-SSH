from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / 'data' / 'curved_real_scale_railway_tunnel_t0t5' / 'raycast_vs_regular_comparison'
SRC = OUT_DIR / 'step6_crown_settlement_workflow_professor_report_with_images.pptx'
OUT = OUT_DIR / 'step6_blender_groundtruth_raycast_regular_validation.pptx'
NOTES = OUT_DIR / 'step6_blender_groundtruth_raycast_regular_validation_notes.md'

BLUE = RGBColor(15, 76, 129)
BLUE2 = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(239, 246, 255)
RED = RGBColor(220, 38, 38)
ORANGE = RGBColor(245, 158, 11)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
GRAY = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(15, 23, 42)

labels = ['T0', 'T1', 'T2', 'T3', 'T4', 'T5']
gt = [0.0, -10.0, -22.0, -38.0, -58.0, -80.0]
regular = [0.0, -9.9, -21.7, -37.6, -57.3, -79.1]
raycast = [0.0, -10.243, -21.632, -37.070, -56.607, -77.900]
regular_err = ['-', '1.00%', '1.36%', '1.05%', '1.21%', '1.13%']
raycast_err = ['-', '2.43%', '1.67%', '2.45%', '2.40%', '2.63%']

prs = Presentation(SRC)
blank = prs.slide_layouts[6]

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

def add_header(slide, title, kicker='VALIDATION LOGIC'):
    add_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    k = slide.shapes.add_textbox(Inches(0.45), Inches(0.11), Inches(3.5), Inches(0.35))
    p = k.text_frame.paragraphs[0]
    p.text = kicker; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHITE
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(12.0), Inches(0.55))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = BLACK

def add_text(slide, x, y, w, h, text, size=14, color=SLATE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    if align:
        p.alignment = align
    return box

def add_box(slide, x, y, w, h, title, body, fill, line):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(1.5)
    add_text(slide, x + 0.18, y + 0.15, w - 0.36, 0.3, title, 14, line, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.75, body, 11.5, SLATE, False, PP_ALIGN.CENTER)
    return shape

def add_arrow(slide, x, y, w=0.7):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.33))
    arr.fill.solid(); arr.fill.fore_color.rgb = GRAY
    arr.line.fill.background()

def add_table(slide, x, y, w, h):
    data = [['Time', 'Ground truth', 'Regular tool', 'Raycast tool', 'Regular error', 'Raycast error']]
    for i, label in enumerate(labels):
        data.append([label, f'{gt[i]:.1f}', f'{regular[i]:.3f}'.rstrip('0').rstrip('.'), f'{raycast[i]:.3f}'.rstrip('0').rstrip('.'), regular_err[i], raycast_err[i]])
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [0.75, 1.35, 1.35, 1.35, 1.2, 1.2]
    for c, width in enumerate(widths):
        table.columns[c].width = Inches(width)
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9.5 if r else 9)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else BLACK
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE if r == 0 else (LIGHT if r % 2 else WHITE)
            cell.margin_left = Inches(0.03); cell.margin_right = Inches(0.03)

s = prs.slides.add_slide(blank)
add_header(s, 'Validation logic: compare three versions of the same tunnel')
add_text(s, 0.75, 1.25, 11.8, 0.38, 'Purpose: prove how much the Step 6 tool deviates from known Blender ground truth.', 15, BLUE, True, PP_ALIGN.CENTER)
add_box(s, 0.55, 2.0, 2.15, 1.35, '3D Blender model', 'Curved railway tunnel with lining, rail, pipes, cables, lights and equipment.', LIGHT_BLUE, BLUE2)
add_arrow(s, 2.85, 2.52)
add_box(s, 3.55, 1.55, 2.05, 1.0, 'Ground truth', 'Known crown settlement: 0, -10, -22, -38, -58, -80 mm.', RGBColor(254, 242, 242), RED)
add_box(s, 3.55, 2.9, 2.05, 1.0, 'Same location', 'Crown / Dinh ham at Ch 52.0m for every epoch.', RGBColor(240, 253, 244), GREEN)
add_arrow(s, 5.78, 2.52)
add_box(s, 6.5, 1.42, 2.1, 1.2, 'Regular clean', 'Direct clean lining export. Tests the ideal algorithm condition.', RGBColor(239, 246, 255), BLUE2)
add_box(s, 6.5, 3.02, 2.1, 1.2, 'Raycast TLS', 'Simulated scanner with noise, occlusion and dropout. Tests field-like condition.', RGBColor(255, 247, 237), ORANGE)
add_arrow(s, 8.8, 2.52)
add_box(s, 9.55, 2.0, 2.15, 1.35, 'Step 6 tool', 'Measure crown settlement for T0-T5 using the same crown point.', RGBColor(240, 253, 244), GREEN)
add_arrow(s, 11.85, 2.52, 0.55)
add_box(s, 12.25, 2.0, 0.85, 1.35, 'Error', 'mm\n%', LIGHT, RED)
add_text(s, 1.0, 5.0, 5.25, 0.55, 'Error mm = Tool result - Ground truth', 17, BLACK, True, PP_ALIGN.CENTER)
add_text(s, 7.05, 5.0, 5.25, 0.55, 'Error % = |Error mm| / |Ground truth| x 100', 17, BLACK, True, PP_ALIGN.CENTER)
add_text(s, 1.1, 6.18, 11.1, 0.38, 'T0 is excluded from MAPE because the ground-truth settlement is 0 mm.', 12, SLATE, False, PP_ALIGN.CENTER)

s = prs.slides.add_slide(blank)
add_header(s, 'Three-way comparison result: ground truth vs regular vs raycast', 'ACCURACY RESULT')
add_table(s, 0.6, 1.45, 7.2, 4.1)
add_box(s, 8.15, 1.55, 1.85, 1.1, 'Regular MAPE', '1.15%', RGBColor(240, 253, 244), GREEN)
add_box(s, 10.35, 1.55, 1.85, 1.1, 'Raycast MAPE', '2.315%', RGBColor(255, 247, 237), ORANGE)
add_text(s, 8.1, 3.25, 4.35, 0.48, 'Interpretation', 18, BLUE, True)
add_text(s, 8.1, 3.82, 4.55, 1.55, 'Regular clean has lower error because it is a clean surface export. Raycast is harder because it simulates field TLS effects such as noise, occlusion and missing points.', 14, SLATE)
add_text(s, 8.1, 5.55, 4.55, 0.78, 'Main conclusion: both branches follow the same settlement trend from T0 to T5, so Step 6 can quantify crown settlement reliably.', 14, GREEN, True)
add_text(s, 0.65, 6.35, 11.9, 0.35, 'Main metric only: Crown settlement / Lun dinh ham at Ch 52.0m. M3C2/p95 are supplementary, not the final conclusion.', 11.5, SLATE, False, PP_ALIGN.CENTER)

prs.save(OUT)

notes = '''# Step 6 Workflow Explanation - Blender -> Ground Truth -> Regular/Raycast -> Tool Error

## One-minute explanation
Dau tien, em dung mot mo hinh ham duong sat cong trong Blender. Vi mo hinh nay do minh kiem soat, em dat truoc do lun that tai dinh ham cho 6 moc thoi gian T0-T5. Day la ground truth.

Tu cung mot mo hinh do, em tao hai bo du lieu kiem thu. Bo thu nhat la regular clean, tuc la xuat truc tiep be mat lining sach de kiem tra thuat toan trong dieu kien ly tuong. Bo thu hai la raycast field-like, tuc la mo phong may quet laser TLS ban tia tu nhieu tram doc theo ham cong, nen du lieu co noise, occlusion va dropout giong thuc dia hon.

Sau do em chay tool Step 6 tren ca hai bo du lieu. Tool luon do cung mot thong so la crown settlement, tuc lun dinh ham, tai cung vi tri Ch 52.0m. Cuoi cung em so ket qua tool voi ground truth de tinh sai lech theo mm va phan tram.

Ket qua: regular clean co MAPE 1.15%, raycast field-like co MAPE 2.315%. Raycast sai so cao hon vi giong thuc dia hon, nhung ca hai van bam dung xu huong lun tu T0 den T5.

## Formula
- Error mm = Tool result - Ground truth
- Error % = |Error mm| / |Ground truth mm| x 100
- T0 khong dung de tinh MAPE vi ground truth bang 0 mm.

## Main result table
| Time | Ground truth mm | Regular tool mm | Raycast tool mm | Regular error % | Raycast error % |
|---|---:|---:|---:|---:|---:|
| T0 | 0.0 | 0.0 | 0.0 | - | - |
| T1 | -10.0 | -9.9 | -10.243 | 1.00% | 2.43% |
| T2 | -22.0 | -21.7 | -21.632 | 1.36% | 1.67% |
| T3 | -38.0 | -37.6 | -37.070 | 1.05% | 2.45% |
| T4 | -58.0 | -57.3 | -56.607 | 1.21% | 2.40% |
| T5 | -80.0 | -79.1 | -77.900 | 1.13% | 2.63% |

## Key message for professor
Tool Step 6 is validated against a known Blender ground truth. The clean branch checks algorithm accuracy, and the raycast branch checks field-like robustness. The main measurement is crown settlement at Ch 52.0m, with average error 1.15% for clean data and 2.315% for raycast TLS data.
'''
NOTES.write_text(notes, encoding='utf-8')
print(OUT)
print(NOTES)
