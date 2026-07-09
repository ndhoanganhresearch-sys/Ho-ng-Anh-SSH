"""Attach English speaker-notes (spoken script) to each slide of the EN deck.
Re-run AFTER rebuilding the deck from build_step6_deck_en.js (rebuild wipes notes).
"""
import sys
from pptx import Presentation

PATH = "Step6_T0_Tn_Deformation_EN.pptx"

NOTES = [
# 1 Title
"""Step 6 is the deformation stage. We compare a baseline scan, T0, with a later scan, Tn, and measure how the tunnel surface moved between them.""",
# 2 Context
"""Step 6 is the sixth of seven steps. It runs after T0 and Tn are cleaned and registered, and does four things: measure displacement, quantify it, classify risk, and optionally forecast.""",
# 3 T0 vs Tn
"""T0 is the undeformed baseline; Tn is the later scan. Rule one: register first, or the difference is just setup error. Registration must keep local deformation, not erase it.""",
# 4 Workflow
"""The pipeline: 6.1 load both epochs, 6.3 the M3C2 map, then per-section parameters, then warnings. That is the full two-scan job. More epochs add the trend and forecast.""",
# 5 M3C2 method
"""M3C2: at each point we project the T0-to-Tn change onto the surface normal, giving a signed distance. Negative means inward. Only values above the LoD, the noise limit, count as real.""",
# 6 Real result
"""Real data, not a mock-up. Both views show three damage zones: crown at twenty metres, sidewall at forty-five, local damage at sixty-five, matching ground truth. Peak: minus forty-four millimetres.""",
# 7 Parameters
"""Four parameters per section: crown settlement, lateral convergence, ovality and eccentricity. Each has a caution and a critical threshold that drive the colour-coding.""",
# 8 Risk classification
"""Sections are flagged versus T0. The local gate stops a uniform bias painting the whole tunnel; ovality and eccentricity must also be local anomalies. Dimension changes use the absolute threshold. One classifier feeds every view.""",
# 9 Real example
"""On the complex dataset: thirty-eight critical, nine caution, thirty-three OK. Crown of ninety-two millimetres versus ninety ground truth, and one hundred percent recall on the damage band.""",
# 10 Extension
"""With three or more epochs you get trend and forecast. Track the peak, not the median, which stays near zero. The forecast predicts time-to-threshold; trust it only when R-squared is high.""",
# 11 Meaning & cautions
"""Reading rules: negative is inward, always read with chainage; ignore values below the LoD; watch the coverage warning; registration keeps deformation; prefer the local peak; trust forecasts only with a good fit.""",
# 12 Summary
"""In one sentence: Step 6 compares Tn to T0, measures displacement with M3C2, quantifies deformation per section, and flags caution or critical by chainage. Thank you.""",
]

p = Presentation(PATH)
assert len(NOTES) == len(p.slides), f"notes {len(NOTES)} != slides {len(p.slides)}"
for sl, txt in zip(p.slides, NOTES):
    sl.notes_slide.notes_text_frame.text = txt.strip()
try:
    p.save(PATH)
    out = PATH
except PermissionError:
    out = PATH.replace(".pptx", "_NOTES.pptx")
    p.save(out)
    print(f"[warn] {PATH} is locked (open in PowerPoint?). Saved to {out} instead.")
print(f"Attached {len(NOTES)} speaker-note scripts to {out}")
